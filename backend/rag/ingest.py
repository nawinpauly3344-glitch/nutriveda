"""
MHB Document Ingestion Pipeline
Reads PDFs and PPTs from the MHB folder, creates embeddings using OpenAI,
stores them in a simple numpy/JSON vector store (Python 3.14 compatible).
"""

import os
import sys
import re
import time
import logging
from pathlib import Path
from typing import List, Dict

# Ensure backend root is in sys.path so 'rag' package is importable
_backend_root = str(Path(__file__).parent.parent)
if _backend_root not in sys.path:
    sys.path.insert(0, _backend_root)

from dotenv import load_dotenv
load_dotenv(dotenv_path=os.path.join(os.path.dirname(__file__), "../../.env"))

log = logging.getLogger(__name__)
logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")

MHB_PATH = os.getenv("MHB_FILES_PATH", r"C:\Users\Deepu\Documents\MHB Files\MHB")


def extract_text_from_pdf(file_path: str) -> str:
    from pypdf import PdfReader
    try:
        reader = PdfReader(file_path)
        texts = []
        for page in reader.pages:
            text = page.extract_text()
            if text and text.strip():
                texts.append(text.strip())
        return "\n\n".join(texts)
    except Exception as e:
        log.warning(f"Failed to read PDF {file_path}: {e}")
        return ""


def extract_text_from_pptx(file_path: str) -> str:
    from pptx import Presentation
    try:
        prs = Presentation(file_path)
        texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            if slide_texts:
                texts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
        return "\n\n".join(texts)
    except Exception as e:
        log.warning(f"Failed to read PPTX {file_path}: {e}")
        return ""


def chunk_text(text: str, chunk_size: int = 1000, overlap: int = 200) -> List[str]:
    """Split text into overlapping chunks."""
    chunks = []
    start = 0
    text = text.strip()
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap
        if start >= len(text):
            break
    return chunks


def load_and_chunk_documents(folder_path: str) -> List[Dict]:
    """Load all PDFs and PPTs and return chunked documents with metadata."""
    folder = Path(folder_path)
    if not folder.exists():
        log.error(f"MHB folder not found: {folder_path}")
        return []

    files = list(folder.glob("*.pdf")) + list(folder.glob("*.pptx")) + list(folder.glob("*.ppt"))
    log.info(f"Found {len(files)} files in {folder_path}")

    all_chunks = []
    for file_path in files:
        file_name = file_path.name
        log.info(f"Processing: {file_name}")

        if file_path.suffix.lower() == ".pdf":
            text = extract_text_from_pdf(str(file_path))
        elif file_path.suffix.lower() in (".pptx", ".ppt"):
            text = extract_text_from_pptx(str(file_path))
        else:
            continue

        if not text.strip():
            log.warning(f"  No text extracted from: {file_name}")
            continue

        # Clean up topic name from file name
        topic = file_path.stem.replace("_", " ").replace("-", " ").strip()
        topic = re.sub(r"^\d+\s*-?\s*\d*\s*-?\s*", "", topic).strip()

        chunks = chunk_text(text, chunk_size=1000, overlap=200)
        for i, chunk in enumerate(chunks):
            all_chunks.append({
                "text": chunk,
                "source": file_name,
                "topic": topic,
                "file_type": file_path.suffix.lower(),
                "chunk_index": i,
                "total_chunks": len(chunks),
            })

        log.info(f"  ✓ {file_name}: {len(text)} chars → {len(chunks)} chunks")

    log.info(f"Total: {len(all_chunks)} chunks from {len(files)} files")
    return all_chunks


def run_ingestion(force_rebuild: bool = False):
    """Main ingestion pipeline. Call this once to build the knowledge base."""
    from rag.vectorstore import vectorstore_exists, embed_texts, save_vectorstore

    log.info("=" * 60)
    log.info("  MHB Nutrition Knowledge Base - Ingestion Pipeline")
    log.info("=" * 60)

    # Check if already built
    if not force_rebuild:
        exists, count = vectorstore_exists()
        if exists:
            log.info(f"Vector store already exists: {count} chunks")
            log.info("Use --force flag to rebuild from scratch")
            return

    # Load and chunk all documents
    chunks = load_and_chunk_documents(MHB_PATH)
    if not chunks:
        log.error("No documents loaded. Check MHB_FILES_PATH in .env")
        return

    # Extract texts for embedding
    texts = [c["text"] for c in chunks]
    log.info(f"\nCreating embeddings for {len(texts)} chunks using OpenAI text-embedding-3-small...")
    log.info("(This may take a few minutes and costs ~$0.01-0.05)")

    # Embed in batches with rate limit pauses
    import numpy as np
    from rag.vectorstore import embed_texts, save_vectorstore

    all_embeddings = []
    batch_size = 50
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i + batch_size]
        log.info(f"Embedding batch {i//batch_size + 1}/{(len(texts)-1)//batch_size + 1}...")
        embeddings = embed_texts(batch)
        all_embeddings.append(embeddings)
        if i + batch_size < len(texts):
            time.sleep(0.3)

    all_embeddings_np = np.vstack(all_embeddings)

    # Save vector store
    save_vectorstore(all_embeddings_np, chunks)

    log.info("=" * 60)
    log.info(f"✓ Knowledge base ready: {len(chunks)} chunks indexed")
    log.info("=" * 60)


if __name__ == "__main__":
    force = "--force" in sys.argv
    run_ingestion(force_rebuild=force)
